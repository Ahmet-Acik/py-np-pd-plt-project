from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Define slide layout
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]
title_only_layout = prs.slide_layouts[5]

# Helper function to add a slide with title and content
def add_slide(title, content_lines):
    slide = prs.slides.add_slide(content_slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = "\n".join(content_lines)
    return slide

# Title Slide
slide_0 = prs.slides.add_slide(title_slide_layout)
slide_0.shapes.title.text = "How do computers talk to each other?"
slide_0.placeholders[1].text = "Exploring Peer-to-Peer and Client-Server Networks"

# Slide 1 - Hook
add_slide("Hook: Streaming vs Sharing", [
    "Imagine one illegal movie download being shared 10 million times — no website involved!",
    "Or your school login not working because a server is down.",
    "Sharing music via AirDrop vs. streaming from Spotify — what's the difference?"
])

# Slide 2 - Peer-to-Peer Networks
add_slide("Peer-to-Peer (P2P) – Like Bluetooth Sharing", [
    "No central server or admin",
    "Devices connect directly to share files",
    "Each device acts as both sender and receiver",
    "Examples: Bluetooth, AirDrop, Minecraft LAN, file-sharing apps"
])

# Slide 3 - Client-Server Networks
add_slide("Client-Server – Like Streaming or Logging In", [
    "One main server provides services to others (clients)",
    "Devices request files, login access, or websites",
    "Easier to control and manage",
    "Examples: School logins, Google Classroom, Fortnite online"
])

# Slide 4 - Comparing the Two Models
add_slide("P2P vs Client-Server – Spot the Differences", [
    "P2P: Shared control, low cost, hard to secure, e.g., Minecraft LAN",
    "Client-Server: Central control, higher cost, secure, e.g., Google Classroom"
])

# Slide 5 - Quick Quiz
add_slide("Quick Quiz - Which Model?", [
    "1. Sam shares music via AirDrop → ?",
    "2. Jade logs into school email → ?",
    "3. Tom and Mia play Minecraft in same room → ?",
    "Hold up 1 finger for P2P, 2 for Client-Server!"
])

# Slide 6 - Why It Matters
add_slide("Why does this matter in real life?", [
    "Helps choose the right network",
    "Schools need security → Client-Server",
    "No internet? Use Peer-to-Peer with friends"
])

# Save the presentation
pptx_path = "Peer_to_Peer_vs_Client_Server_Lesson.pptx"
prs.save(pptx_path)

# pptx_path
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

# Reload the presentation
prs = Presentation(pptx_path)

# Function to add a basic icon (shape) to a slide
def add_icon(slide, left_inch, top_inch, icon_text):
    left = Inches(left_inch)
    top = Inches(top_inch)
    width = Inches(1)
    height = Inches(1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.text = icon_text
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(91, 155, 213)  # Light blue
    text_frame = shape.text_frame
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True

# Add icons/visuals to each slide except title
icon_labels = [
    "🎧",  # Hook
    "📲",  # P2P
    "🌐",  # Client-Server
    "⚖️",  # Comparison
    "❓",  # Quiz
    "🔐"   # Recap
]

for i, icon in enumerate(icon_labels, start=1):
    slide = prs.slides[i]
    add_icon(slide, 6.5, 0.5, icon)

# Save the updated presentation
updated_pptx_path = "Peer_to_Peer_vs_Client_Server_Lesson_Icons.pptx"
prs.save(updated_pptx_path)


